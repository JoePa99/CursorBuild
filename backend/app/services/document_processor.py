import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import pandas as pd
from typing import List, Dict, Any, Optional
import logging
from pathlib import Path

from app.models.document import Document, DocumentType, DocumentStatus, DocumentChunk
from app.core.config import settings

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Document processing service for text extraction and chunking"""
    
    def __init__(self):
        self.supported_types = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.txt': self._extract_txt,
            '.csv': self._extract_csv,
            '.xlsx': self._extract_excel,
            '.pptx': self._extract_pptx
        }
    
    async def process_document(self, document: Document) -> Document:
        """Process a document to extract text and create chunks"""
        try:
            logger.info(f"Processing document: {document.filename}")
            
            # Update status to processing
            document.status = DocumentStatus.PROCESSING
            
            # Extract text based on document type
            if document.document_type.value in self.supported_types:
                extractor = self.supported_types[document.document_type.value]
                extracted_text = await extractor(document.file_path)
            else:
                raise ValueError(f"Unsupported document type: {document.document_type}")
            
            # Clean and normalize text
            cleaned_text = self._clean_text(extracted_text)
            document.extracted_text = cleaned_text
            
            # Create text chunks
            chunks = self._create_chunks(cleaned_text, document.id)
            document.chunk_count = len(chunks)
            
            # Update status to processed
            document.status = DocumentStatus.PROCESSED
            
            logger.info(f"Successfully processed document: {document.filename}")
            return document, chunks
            
        except Exception as e:
            logger.error(f"Error processing document {document.filename}: {str(e)}")
            document.status = DocumentStatus.FAILED
            document.processing_errors.append(str(e))
            raise
    
    async def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyMuPDF"""
        try:
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            raise
    
    async def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX using python-docx"""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            raise
    
    async def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error extracting TXT text: {str(e)}")
            raise
    
    async def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"Error extracting CSV text: {str(e)}")
            raise
    
    async def _extract_excel(self, file_path: str) -> str:
        """Extract text from Excel file"""
        try:
            df = pd.read_excel(file_path)
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"Error extracting Excel text: {str(e)}")
            raise
    
    async def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            raise
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Remove common OCR artifacts
        text = text.replace('\x00', '')  # Null bytes
        text = text.replace('\r', '\n')  # Carriage returns
        
        # Normalize line breaks
        text = text.replace('\n\n\n', '\n\n')
        
        # Remove special characters that might cause issues
        text = text.replace('\t', ' ')
        
        return text.strip()
    
    def _create_chunks(self, text: str, document_id: str, chunk_size: int = 1000, overlap: int = 200) -> List[DocumentChunk]:
        """Create semantically coherent text chunks"""
        if not text:
            return []
        
        chunks = []
        words = text.split()
        current_chunk = []
        current_length = 0
        chunk_index = 0
        
        for word in words:
            current_chunk.append(word)
            current_length += len(word) + 1  # +1 for space
            
            if current_length >= chunk_size:
                # Create chunk
                chunk_text = ' '.join(current_chunk)
                chunk = DocumentChunk(
                    document_id=str(document_id),
                    chunk_index=chunk_index,
                    content=chunk_text,
                    start_position=len(' '.join(words[:len(words) - len(current_chunk)])),
                    end_position=len(' '.join(words[:len(words) - len(current_chunk) + len(current_chunk)]))
                )
                chunks.append(chunk)
                
                # Prepare for next chunk with overlap
                overlap_words = current_chunk[-overlap//10:] if overlap > 0 else []
                current_chunk = overlap_words
                current_length = sum(len(word) + 1 for word in overlap_words)
                chunk_index += 1
        
        # Add final chunk if there's remaining text
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            chunk = DocumentChunk(
                document_id=str(document_id),
                chunk_index=chunk_index,
                content=chunk_text,
                start_position=len(text) - len(chunk_text),
                end_position=len(text)
            )
            chunks.append(chunk)
        
        return chunks
    
    def get_document_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from document"""
        metadata = {}
        file_path = Path(file_path)
        
        if file_path.exists():
            metadata['file_size'] = file_path.stat().st_size
            metadata['last_modified'] = file_path.stat().st_mtime
            
            # Extract additional metadata based on file type
            if file_path.suffix.lower() == '.pdf':
                try:
                    doc = fitz.open(str(file_path))
                    metadata['page_count'] = len(doc)
                    metadata['title'] = doc.metadata.get('title', '')
                    metadata['author'] = doc.metadata.get('author', '')
                    doc.close()
                except Exception as e:
                    logger.warning(f"Could not extract PDF metadata: {e}")
        
        return metadata 